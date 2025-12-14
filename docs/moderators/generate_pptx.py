"""
Generate a PowerPoint deck for the workshop.

Output: docs/moderators/FOLIENSKRIPT_WORKSHOP.pptx

Requires:
  python3 -m pip install --user python-pptx
"""

from __future__ import annotations

from datetime import date
from pathlib import Path
from typing import Iterable, Optional, Sequence


def _ensure_dep():
    try:
        import pptx  # noqa: F401
    except Exception as e:  # pragma: no cover
        raise SystemExit(
            "Missing dependency 'python-pptx'. Install with:\n"
            "  python3 -m pip install --user python-pptx\n"
            f"\nOriginal error: {e}"
        )


def _set_notes(slide, text: str) -> None:
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    tf.text = (text or "").strip()


def _add_title(slide, title: str, subtitle: Optional[str] = None) -> None:
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    left = Inches(0.8)
    top = Inches(1.2)
    width = Inches(11.8)
    height = Inches(1.2)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        tb2 = slide.shapes.add_textbox(left, top + Inches(1.0), width, Inches(0.8))
        tf2 = tb2.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(20)
        p2.alignment = PP_ALIGN.LEFT


def _add_header(slide, title: str) -> None:
    from pptx.util import Inches, Pt

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.8), Inches(0.6))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(30)
    run.font.bold = True


def _add_bullets(slide, bullets: Sequence[str], top: float = 1.4) -> None:
    from pptx.util import Inches, Pt

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.6), Inches(5.6))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True

    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = b
        p.level = 0
        p.font.size = Pt(20)


def _add_code_anchor(slide, title: str, items: Sequence[str], top: float = 5.9) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.6), Inches(1.2))
    tf = tb.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    r0.font.size = Pt(14)
    r0.font.bold = True
    r0.font.color.rgb = RGBColor(90, 90, 90)

    for it in items:
        p = tf.add_paragraph()
        p.text = it
        p.level = 1
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(90, 90, 90)


def _add_pipeline_diagram(slide) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN

    # 4 rounded rectangles with arrows
    labels = ["Reader", "Summarizer", "Critic", "Integrator"]
    x0 = Inches(0.9)
    y = Inches(2.6)
    w = Inches(2.8)
    h = Inches(1.0)
    gap = Inches(0.5)

    fill = RGBColor(245, 248, 255)
    line = RGBColor(90, 120, 180)

    boxes = []
    for i, lab in enumerate(labels):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x0 + i * (w + gap),
            y,
            w,
            h,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = line
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = lab
        p.font.size = Pt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        boxes.append(shape)

    # arrows
    for i in range(len(boxes) - 1):
        a = boxes[i]
        b = boxes[i + 1]
        x1 = a.left + a.width
        y1 = a.top + a.height // 2
        x2 = b.left
        y2 = b.top + b.height // 2
        conn = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
        conn.line.color.rgb = line
        conn.line.width = Pt(2)


def _slide(prs, kind: str = "content"):
    if kind == "title":
        return prs.slides.add_slide(prs.slide_layouts[6])  # blank
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank for consistent layout


def build_deck(out_path: Path) -> None:
    _ensure_dep()
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    today = date.today().isoformat()

    slides = [
        {
            "title": "Multi-Agent Orchestration Workshop",
            "subtitle": "Paper Analyzer • LangChain vs. LangGraph vs. DSPy",
            "bullets": [],
            "notes": (
                "Ziel: Multi-Agent Orchestration anhand eines konkreten Projekts verstehen.\n"
                "Wir gehen von Theorie → Prinzip → Code → Aufgaben.\n"
                f"Deck generiert am {today}."
            ),
            "type": "title",
        },
        {
            "title": "Ablauf (60 Minuten)",
            "bullets": [
                "Einleitung + Mini-Demo (5 min)",
                "Praxis: LangChain → LangGraph → DSPy (50 min)",
                "Wrap-Up: Erkenntnisse + Transfer + Q&A (5 min)",
            ],
            "notes": (
                "Timing-Ansage: Wir beenden alle Aktivitäten ca. 5 Minuten vor Schluss.\n"
                "Wenn Setup hängt: auf Code-Walkthrough ohne Ausführen wechseln."
            ),
        },
        {
            "title": "Lernziele",
            "bullets": [
                "Multi-Agent Prinzip: Zerlegen → Orchestrieren → Bewerten → Integrieren",
                "Im Code wiederfinden: Agenten (`app/agents/*`) + Workflows (`app/workflows/*`)",
                "Zwei Aufgabenarten: (1) Verständnis, (2) Grenzen/Trade-offs",
            ],
            "notes": (
                "Regel für den Workshop: Erst Prinzip erklären, dann zeigen wo es im Code steckt, dann Aufgaben."
            ),
        },
        {
            "title": "Pipeline als mentale Landkarte",
            "bullets": [
                "Reader: extrahiert strukturierte Notes (Ground Truth im Workflow)",
                "Summarizer: schreibt Summary aus Notes",
                "Critic: bewertet Summary gegen Notes (Rubrik 0–5)",
                "Integrator: Meta Summary aus Summary + Critic (weiterhin grounded in Notes)",
            ],
            "notes": (
                "Wichtig: Notes sind die Referenz. Critic bestraft ungrounded claims.\n"
                "Integrator nutzt Critic-Signal, erfindet aber keine Fakten."
            ),
            "diagram": "pipeline",
            "anchors": [
                "Reader: `app/agents/reader.py`",
                "Summarizer: `app/agents/summarizer.py`",
                "Critic: `app/agents/critic.py`",
                "Integrator: `app/agents/integrator.py`",
            ],
        },
        {
            "title": "Live-Demo (90 Sekunden)",
            "bullets": [
                "App öffnen → Tab „Analysis“",
                "Kurzes PDF/TXT hochladen",
                "Pipeline: LangChain auswählen",
                "„Analyze“ → Artefakte + Timing anschauen",
            ],
            "notes": (
                "Zeige: Meta Summary, Summary, Notes, Critic, Timing.\n"
                "Plenum-Frage: Welcher Schritt ist am teuersten (Reader vs Summarizer) und warum?"
            ),
            "anchors": [
                "UI: `app/app.py`",
                "Telemetry: `app/telemetry.py`",
            ],
        },
        {
            "title": "Praxis-Regeln",
            "bullets": [
                "Immer nur 1 Änderung → erneut ausführen → Effekt beobachten",
                "Hypothese → Änderung → Messung (Output + Timing/Telemetry)",
                "Bei Problemen: „Debug Mode“ in der Sidebar aktivieren",
            ],
            "notes": (
                "Erwartungsmanagement: Wir optimieren nicht perfekt, sondern lernen Mechanismen & Trade-offs."
            ),
        },
        {
            "title": "LangChain – Prinzip",
            "bullets": [
                "Linearer Kontrollfluss: Step A liefert Input für Step B",
                "Abhängigkeiten sind implizit (durch Reihenfolge im Code)",
                "Sehr gut für Einstieg & schnelle Pipelines",
            ],
            "notes": (
                "Hauptpunkt: Reihenfolge = Kontrollfluss. Keine expliziten Branches/Loops."
            ),
            "anchors": [
                "`app/workflows/langchain_pipeline.py` (`run_pipeline`) ",
            ],
        },
        {
            "title": "LangChain – Wo ist das im Code?",
            "bullets": [
                "Preprocessing: `build_analysis_context(...)`",
                "Reader → Summarizer → Critic → Integrator (sequenziell)",
                "Timing wird pro Stage gemessen und als Telemetrie geloggt",
            ],
            "notes": (
                "Kurz durch `run_pipeline()` scrollen: Aufruf-Reihenfolge + welche Daten weitergereicht werden."
            ),
            "anchors": [
                "`app/utils.py` (`build_analysis_context`) ",
                "`app/workflows/langchain_pipeline.py`",
                "`app/telemetry.py` (`log_row`) ",
            ],
        },
        {
            "title": "Aufgabe (Verständnis): Prompt-Wirkung",
            "bullets": [
                "In `app/agents/summarizer.py`: „200-300 words“ → „50-100 words“",
                "LangChain erneut ausführen",
                "Beobachten: Summary-Länge + Critic (Coverage/Specificity) + Meta Summary",
            ],
            "notes": (
                "Diskussion: Welche Nebenwirkung hat ‚kürzer‘ auf Coverage/Specificity?\n"
                "Grenze: Wenn Notes keine Zahlen enthalten, muss Summary ‚not reported‘ sagen."
            ),
        },
        {
            "title": "Aufgabe (Grenzen): Reihenfolge/Abhängigkeiten",
            "bullets": [
                "Prinzip: Critic braucht Notes + Summary",
                "In `app/workflows/langchain_pipeline.py`: Critic vor Summarizer ziehen",
                "Erwartung: Bewertung wird sinnlos/konzeptionell kaputt",
            ],
            "notes": (
                "Plenum-Frage: Wie macht man Abhängigkeiten sichtbar, ohne gleich LangGraph zu nutzen?"
            ),
        },
        {
            "title": "LangGraph – Prinzip",
            "bullets": [
                "Explizite Nodes + Edges + gemeinsamer State",
                "Conditional Flow: Branches & Loops möglich",
                "Visualisierung als Graph (DOT/Graphviz) hilft beim Debugging",
            ],
            "notes": (
                "Hauptpunkt: Kontrollfluss ist explizit und kann anhand von Scores/Heuristiken geroutet werden."
            ),
            "anchors": [
                "`app/workflows/langgraph_pipeline.py`",
            ],
        },
        {
            "title": "LangGraph – Wo ist das im Code?",
            "bullets": [
                "`PipelineState`: welche Felder die Nodes lesen/schreiben",
                "`_build_langgraph_workflow()`: Nodes/Edges definieren",
                "`_critic_post_path()`: Routing-Entscheidungen (z. B. rework-loop)",
                "Output enthält zusätzlich `graph_dot` + ggf. `critic_loops`",
            ],
            "notes": (
                "Zeige kurz: graph.add_node/add_edge/add_conditional_edges.\n"
                "Dann in der App den Graph anzeigen."
            ),
            "anchors": [
                "`app/workflows/langgraph_pipeline.py` (`PipelineState`) ",
                "`app/workflows/langgraph_pipeline.py` (`_build_langgraph_workflow`) ",
                "`app/workflows/langgraph_pipeline.py` (`_critic_post_path`) ",
            ],
        },
        {
            "title": "Aufgabe (Verständnis): State & Zusatz-Nodes",
            "bullets": [
                "`_execute_translator_node`: setzt `state['summary_translated']`",
                "`_execute_keyword_node`: setzt `state['keywords']`",
                "Ändere eine Default-Option (z. B. EN + ultra_short) und führe LangGraph erneut aus",
            ],
            "notes": (
                "Variante A: In `_execute_translator_node` Defaults ändern.\n"
                "Variante B: (optional) `translator_language`/`translator_style` in `app/app.py` ins config-Dict aufnehmen."
            ),
            "anchors": [
                "`app/workflows/langgraph_pipeline.py` (`_execute_translator_node`) ",
                "`app/workflows/langgraph_pipeline.py` (`_execute_keyword_node`) ",
            ],
        },
        {
            "title": "Aufgabe (Grenzen): Routing/Qualitätsschwellen",
            "bullets": [
                "Finde Routing-Logik in `_critic_post_path()`",
                "Ändere eine Schwelle (z. B. strenger/lockerer für rework-loop)",
                "Beobachte: `critic_loops`, Laufzeit, Output-Qualität",
            ],
            "notes": (
                "Diskussion: Zu aggressives Loopen erhöht Kosten/Latenz und kann overfitten.\n"
                "Frage: Welche Signale taugen für Routing (LLM-Score vs Heuristik-F1 vs Kombination)?"
            ),
        },
        {
            "title": "DSPy – Prinzip",
            "bullets": [
                "Deklarativ: `dspy.Signature` beschreibt Inputs/Outputs + Constraints",
                "DSPy generiert Prompts aus Signatures (nicht als Prompt-Strings im Code)",
                "Optional: Teleprompting optimiert anhand eines Dev-Sets",
            ],
            "notes": (
                "Hauptpunkt: Kontrolle verschiebt sich – weniger Prompt-Handarbeit, mehr Zielbeschreibung + Daten."
            ),
            "anchors": [
                "`app/workflows/dspy_pipeline.py`",
            ],
        },
        {
            "title": "DSPy – Wo ist das im Code?",
            "bullets": [
                "Signatures: `ReadNotes`, `Summarize`, `Critique`, `Integrate`",
                "`PaperPipeline.forward(...)`: ruft die Module in Reihenfolge auf",
                "Teleprompting: `_teleprompt_if_requested(...)` (BootstrapFewShot)",
            ],
            "notes": (
                "Wenn DSPy nicht installiert ist: Stub-Modus (App crasht nicht), Konzept trotzdem erklärbar."
            ),
            "anchors": [
                "`app/workflows/dspy_pipeline.py` (Signatures) ",
                "`app/workflows/dspy_pipeline.py` (`_teleprompt_if_requested`) ",
            ],
        },
        {
            "title": "Aufgabe (Verständnis): Signature ändern",
            "bullets": [
                "In `Summarize`-Docstring Fokus ändern (z. B. „Results first, numeric metrics“) ",
                "DSPy ausführen und Output vergleichen",
                "Beobachtung: Verhalten ändert sich ohne Prompt-String-Edit",
            ],
            "notes": (
                "Explizit machen: In LangChain/LangGraph editiert ihr Prompts direkt; DSPy verändert Verhalten über Signatures."
            ),
        },
        {
            "title": "Aufgabe (Grenzen): Teleprompting Trade-off",
            "bullets": [
                "Teleprompting aktivieren (wenn Dev-Set vorhanden)",
                "„Run Teleprompt Comparison“ starten",
                "Vergleichen: F1 Gain vs. Latenz",
            ],
            "notes": (
                "Diskussion: Schlechte/biasierte Dev-Beispiele können Teleprompting in die falsche Richtung treiben.\n"
                "Frage: Wie sähe ein gutes Dev-Set für eure Domäne aus?"
            ),
        },
        {
            "title": "Wrap-Up: 3 Paradigmen auf 1 Slide",
            "bullets": [
                "LangChain: schnell & simpel – aber Abhängigkeiten implizit",
                "LangGraph: expliziter Flow + Routing + Visualisierung – mehr Boilerplate, dafür Kontrolle",
                "DSPy: deklarativ + (optional) self-improving – braucht Dev-Set, anderes Debugging",
            ],
            "notes": (
                "Plenum-Frage: Für welchen Use Case würdet ihr heute welches Paradigma wählen – und warum?"
            ),
        },
        {
            "title": "Wrap-Up: Zentrale Erkenntnisse",
            "bullets": [
                "Artefakte pro Schritt (Notes/Summary/Critic/Meta) machen Debugging leichter",
                "Qualität kommt durch Grounding + Feedback + sinnvollen Kontrollfluss",
                "Messbarkeit: Timing/Telemetry + (optional) Quality-Heuristiken",
            ],
            "notes": (
                "Code-Anchor: `app/utils.py` (Robustheit), `app/telemetry.py` (Messung)."
            ),
            "anchors": [
                "`app/utils.py`",
                "`app/telemetry.py`",
            ],
        },
        {
            "title": "Wrap-Up: Transferfragen",
            "bullets": [
                "Wenn wir einen 5. Agenten hinzufügen (z. B. FactChecker): Wo platzieren – und warum?",
                "Wie verhindern wir Halluzinationen: Reader-Prompt, Critic-Score, Routing, Parser?",
                "Was passiert bei sehr großem Input: Token-Limits, Zeit, Kosten – und Gegenmaßnahmen?",
            ],
            "notes": (
                "Abschluss-Satz: Erst Artefakte + Messung schaffen, dann optimieren.\n"
                "Letzte Q&A-Runde."
            ),
        },
    ]

    for s in slides:
        slide = _slide(prs, kind=s.get("type", "content"))
        if s.get("type") == "title":
            _add_title(slide, s["title"], s.get("subtitle"))
        else:
            _add_header(slide, s["title"])
            if s.get("bullets"):
                _add_bullets(slide, s["bullets"])
        if s.get("diagram") == "pipeline":
            _add_pipeline_diagram(slide)
        anchors = s.get("anchors") or []
        if anchors:
            _add_code_anchor(slide, "Code-Anker", anchors, top=6.25)
        _set_notes(slide, s.get("notes", ""))

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> None:
    out = Path(__file__).resolve().parent / "FOLIENSKRIPT_WORKSHOP.pptx"
    build_deck(out)
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()

